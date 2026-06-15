import os
import sys
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from app.config import POPPLER_PATH, DEFAULT_IMAGE_DPI, DEFAULT_IMAGE_QUALITY
from app.utils.file_utils import build_output_path, build_output_path_multi, ensure_output_dir


def pdf_to_docx(input_path: str, output_path: str) -> str:
    # First try Word COM (best quality, needs MS Word installed)
    try:
        import comtypes.client
        word = comtypes.client.CreateObject("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(os.path.abspath(input_path), ConfirmConversions=False)
        doc.SaveAs(os.path.abspath(output_path), FileFormat=16)
        doc.Close()
        word.Quit()
        return output_path
    except Exception:
        pass

    # Fallback — pdf2docx (no Word needed)
    from pdf2docx import Converter as DocxConverter
    cv = DocxConverter(input_path)
    cv.convert(output_path, start=0, end=None, multi_processing=False)
    cv.close()
    return output_path


def pdf_to_pptx(input_path: str, output_path: str) -> str:
    images = convert_from_path(
        input_path,
        dpi=DEFAULT_IMAGE_DPI,
        poppler_path=POPPLER_PATH
    )
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for img in images:
        slide = prs.slides.add_slide(blank_layout)
        tmp_path = output_path + "_tmp.png"
        img.save(tmp_path, format="PNG")
        slide.shapes.add_picture(tmp_path, Inches(0), Inches(0), Inches(10), Inches(7.5))
        os.remove(tmp_path)

    prs.save(output_path)
    return output_path


def pdf_to_images(input_path: str, output_dir: str, target_format: str) -> list[str]:
    images = convert_from_path(
        input_path,
        dpi=DEFAULT_IMAGE_DPI,
        poppler_path=POPPLER_PATH
    )
    ensure_output_dir(output_dir)
    saved = []
    fmt = target_format.upper()
    pil_format = "JPEG" if fmt in ("JPG", "JPEG") else fmt

    for i, img in enumerate(images, start=1):
        out_path = build_output_path_multi(input_path, target_format, output_dir, i)
        if pil_format == "JPEG":
            img = img.convert("RGB")
        img.save(out_path, format=pil_format, quality=DEFAULT_IMAGE_QUALITY)
        saved.append(out_path)

    return saved