import os
import sys
import subprocess
import comtypes.client
from app.utils.file_utils import build_output_path, ensure_output_dir, build_output_path_multi
from app.config import DEFAULT_IMAGE_DPI, DEFAULT_IMAGE_QUALITY, POPPLER_PATH
from pdf2image import convert_from_path
from PIL import Image

if getattr(sys, 'frozen', False):
    os.environ['COMTYPES_CACHE_DIR'] = os.path.join(
        os.path.expanduser("~"), "AppData", "Local", "MorphX", "comtypes_cache"
    )
    comtypes.client.gen_dir = os.environ['COMTYPES_CACHE_DIR']
    os.makedirs(os.environ['COMTYPES_CACHE_DIR'], exist_ok=True)


def _abs(path: str) -> str:
    return os.path.abspath(path)


def _kill_stray_office_process(exe_name: str):
    try:
        subprocess.run(
            ["taskkill", "/F", "/IM", exe_name],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            timeout=5
        )
    except Exception:
        pass


# ── DOCX conversions ────────────────────────────────────────────────────────

def docx_to_pdf(input_path: str, output_path: str) -> str:
    _kill_stray_office_process("WINWORD.EXE")
    word = comtypes.client.CreateObject("Word.Application")
    word.Visible = False
    try:
        doc = word.Documents.Open(_abs(input_path))
        doc.SaveAs(_abs(output_path), FileFormat=17)
        doc.Close()
    finally:
        word.Quit()
    return output_path


def docx_to_pptx(input_path: str, output_path: str) -> str:
    tmp_pdf = _abs(input_path).replace(".docx", "_tmp.pdf")
    docx_to_pdf(input_path, tmp_pdf)
    try:
        result = _pdf_to_pptx_internal(tmp_pdf, output_path)
    finally:
        if os.path.exists(tmp_pdf):
            os.remove(tmp_pdf)
    return result


def docx_to_images(input_path: str, output_dir: str, target_format: str) -> list[str]:
    ensure_output_dir(output_dir)
    tmp_pdf = _abs(input_path).replace(".docx", "_tmp.pdf")
    docx_to_pdf(input_path, tmp_pdf)
    try:
        result = _pdf_to_images_internal(tmp_pdf, input_path, output_dir, target_format)
    finally:
        if os.path.exists(tmp_pdf):
            os.remove(tmp_pdf)
    return result


# ── PPTX / PPT conversions ──────────────────────────────────────────────────

def _open_presentation_via_com(input_path: str, output_path: str, save_format: int = 32) -> str:
    """
    Minimal, proven call pattern (matches the verified working pattern
    used across Stack Overflow / PyPI pptxtopdf / Medium references).
    No extra positional args — Presentations.Open(path) only.
    """
    _kill_stray_office_process("POWERPNT.EXE")
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    try:
        deck = powerpoint.Presentations.Open(_abs(input_path))
        deck.SaveAs(_abs(output_path), save_format)
        deck.Close()
    finally:
        powerpoint.Quit()
    return output_path


def pptx_to_pdf(input_path: str, output_path: str) -> str:
    """
    Requires Microsoft PowerPoint installed (COM automation has no
    cross-platform or Office-free equivalent). If COM fails for any
    reason, falls back to a rasterized rendering via python-pptx +
    Pillow — works for OOXML-based files, not true binary .ppt.
    """
    try:
        return _open_presentation_via_com(input_path, output_path, 32)
    except Exception:
        return _pptx_to_pdf_fallback(input_path, output_path)


def pptx_to_docx(input_path: str, output_path: str) -> str:
    from app.converters.pdf_converter import pdf_to_docx

    tmp_pdf = _abs(input_path).replace(".pptx", "_tmp.pdf")
    pptx_to_pdf(input_path, tmp_pdf)
    try:
        result = pdf_to_docx(tmp_pdf, output_path)
    finally:
        if os.path.exists(tmp_pdf):
            os.remove(tmp_pdf)
    return result


def pptx_to_images(input_path: str, output_dir: str, target_format: str) -> list[str]:
    ensure_output_dir(output_dir)
    tmp_pdf = _abs(input_path).replace(".pptx", "_tmp.pdf")
    pptx_to_pdf(input_path, tmp_pdf)
    try:
        result = _pdf_to_images_internal(tmp_pdf, input_path, output_dir, target_format)
    finally:
        if os.path.exists(tmp_pdf):
            os.remove(tmp_pdf)
    return result


def ppt_to_pdf(input_path: str, output_path: str) -> str:
    """
    Legacy .ppt — same COM call as .pptx (PowerPoint opens both
    transparently). If the file is true pre-2007 binary format and
    COM is unavailable/fails, the fallback will raise a clear
    'Package not found' error since python-pptx cannot parse binary
    .ppt. This is a hard format limitation, not fixable in pure Python.
    """
    try:
        return _open_presentation_via_com(input_path, output_path, 32)
    except Exception:
        return _pptx_to_pdf_fallback(input_path, output_path)


def ppt_to_docx(input_path: str, output_path: str) -> str:
    from app.converters.pdf_converter import pdf_to_docx

    tmp_pdf = _abs(input_path).replace(".ppt", "_tmp.pdf")
    ppt_to_pdf(input_path, tmp_pdf)
    try:
        result = pdf_to_docx(tmp_pdf, output_path)
    finally:
        if os.path.exists(tmp_pdf):
            os.remove(tmp_pdf)
    return result


def ppt_to_images(input_path: str, output_dir: str, target_format: str) -> list[str]:
    ensure_output_dir(output_dir)
    tmp_pdf = _abs(input_path).replace(".ppt", "_tmp.pdf")
    ppt_to_pdf(input_path, tmp_pdf)
    try:
        result = _pdf_to_images_internal(tmp_pdf, input_path, output_dir, target_format)
    finally:
        if os.path.exists(tmp_pdf):
            os.remove(tmp_pdf)
    return result


# ── No-Office fallback renderer ──────────────────────────────────────────────

def _pptx_emu_to_px(emu_value, scale) -> int:
    return int(emu_value * scale) if emu_value else 0


def _pptx_run_font(run, height_px: int):
    from PIL import ImageFont
    size_pt = run.font.size.pt if run.font.size else None
    font_size = int(size_pt * 1.3) if size_pt else max(12, int(height_px * 0.12))
    bold = bool(run.font.bold)
    italic = bool(run.font.italic)
    if bold and italic:
        candidates = ["arialbi.ttf", "arial.ttf"]
    elif bold:
        candidates = ["arialbd.ttf", "arial.ttf"]
    elif italic:
        candidates = ["ariali.ttf", "arial.ttf"]
    else:
        candidates = ["arial.ttf"]
    for name in candidates:
        try:
            return ImageFont.truetype(name, font_size)
        except Exception:
            continue
    return ImageFont.load_default()


def _pptx_run_color(run):
    try:
        color = run.font.color
        if color and color.type is not None and color.rgb is not None:
            rgb = color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return (0, 0, 0)


def _pptx_to_pdf_fallback(input_path: str, output_path: str) -> str:
    import io
    from pptx import Presentation
    from PIL import Image as PILImage, ImageDraw

    prs = Presentation(input_path)
    slide_w_px = 1280
    slide_h_px = int(slide_w_px * prs.slide_height / prs.slide_width)
    scale = slide_w_px / prs.slide_width
    rendered_pages = []

    for slide in prs.slides:
        img = PILImage.new("RGB", (slide_w_px, slide_h_px), color="white")
        draw = ImageDraw.Draw(img)
        for shape in slide.shapes:
            try:
                left = _pptx_emu_to_px(shape.left, scale)
                top = _pptx_emu_to_px(shape.top, scale)
                width = _pptx_emu_to_px(shape.width, scale)
                height = _pptx_emu_to_px(shape.height, scale)

                if shape.shape_type == 13 and shape.image:
                    try:
                        pic_bytes = shape.image.blob
                        pic = PILImage.open(io.BytesIO(pic_bytes)).convert("RGB")
                        pic = pic.resize((max(width, 1), max(height, 1)))
                        img.paste(pic, (left, top))
                        continue
                    except Exception:
                        pass

                if shape.has_text_frame:
                    cursor_y = top + 4
                    for paragraph in shape.text_frame.paragraphs:
                        line_text = "".join(r.text for r in paragraph.runs) or paragraph.text
                        if not line_text.strip():
                            cursor_y += max(14, int(height * 0.12)) if height else 18
                            continue
                        if paragraph.runs:
                            first_run = paragraph.runs[0]
                            font = _pptx_run_font(first_run, height)
                            fill = _pptx_run_color(first_run)
                        else:
                            from PIL import ImageFont
                            font = ImageFont.load_default()
                            fill = (0, 0, 0)
                        draw.text((left + 4, cursor_y), line_text, fill=fill, font=font)
                        try:
                            line_height = font.getbbox(line_text)[3] + 6
                        except Exception:
                            line_height = 20
                        cursor_y += line_height
            except Exception:
                continue
        rendered_pages.append(img)

    if not rendered_pages:
        rendered_pages = [PILImage.new("RGB", (slide_w_px, slide_h_px), color="white")]

    first, rest = rendered_pages[0], rendered_pages[1:]
    first.save(output_path, format="PDF", save_all=True, append_images=rest)
    return output_path


# ── Shared internal helpers ──────────────────────────────────────────────────

def _pdf_to_images_internal(pdf_path, original_input, output_dir, target_format):
    images = convert_from_path(pdf_path, dpi=DEFAULT_IMAGE_DPI, poppler_path=POPPLER_PATH)
    saved = []
    fmt = target_format.upper()
    pil_format = "JPEG" if fmt in ("JPG", "JPEG") else fmt
    for i, img in enumerate(images, start=1):
        out_path = build_output_path_multi(original_input, target_format, output_dir, i)
        if pil_format == "JPEG":
            img = img.convert("RGB")
        img.save(out_path, format=pil_format, quality=DEFAULT_IMAGE_QUALITY)
        saved.append(out_path)
    return saved


def _pdf_to_pptx_internal(pdf_path, output_path):
    from pptx import Presentation
    from pptx.util import Inches
    images = convert_from_path(pdf_path, dpi=DEFAULT_IMAGE_DPI, poppler_path=POPPLER_PATH)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    for i, img in enumerate(images):
        slide = prs.slides.add_slide(blank_layout)
        tmp_img = pdf_path + f"_slide{i}.png"
        img.save(tmp_img, format="PNG")
        slide.shapes.add_picture(tmp_img, Inches(0), Inches(0), Inches(10), Inches(7.5))
        os.remove(tmp_img)
    prs.save(output_path)
    return output_path