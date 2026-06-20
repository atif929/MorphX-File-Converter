import os
from PIL import Image
from app.utils.file_utils import build_output_path, ensure_output_dir
from app.config import DEFAULT_IMAGE_QUALITY
from pptx import Presentation
from pptx.util import Inches

_RGB_ONLY_FORMATS = {"JPEG", "JPG", "PDF", "BMP"}

_PIL_FORMAT_MAP = {
    "JPG":  "JPEG",
    "JPEG": "JPEG",
    "PNG":  "PNG",
    "WEBP": "WEBP",
    "BMP":  "BMP",
    "TIFF": "TIFF",
    "GIF":  "GIF",
    "ICO":  "ICO",
}


def _get_pil_format(fmt: str) -> str:
    return _PIL_FORMAT_MAP.get(fmt.upper(), fmt.upper())


def _prepare_image(img: Image.Image, pil_format: str) -> Image.Image:
    """
    Ensures image mode is compatible with the target format.
    For multi-page images (TIFF), seeks to frame 0 first.
    Converts RGBA/P/LA to RGB for formats that don't support transparency.
    """
    # Handle multi-page — seek to first frame safely
    try:
        img.seek(0)
    except EOFError:
        pass

    if pil_format in ("JPEG", "BMP", "PDF") and img.mode in ("RGBA", "P", "LA"):
        return img.convert("RGB")
    if pil_format == "PNG" and img.mode == "P":
        return img.convert("RGBA")
    return img


def image_to_image(input_path: str, output_path: str, target_format: str) -> str:
    fmt = target_format.upper()
    pil_format = _get_pil_format(fmt)

    with Image.open(input_path) as img:
        img = _prepare_image(img, pil_format)
        if pil_format in ("JPEG", "WEBP"):
            img.save(output_path, format=pil_format, quality=DEFAULT_IMAGE_QUALITY)
        elif pil_format == "TIFF":
            img.save(output_path, format="TIFF", compression="lzw")
        else:
            img.save(output_path, format=pil_format)

    return output_path


# New function to handle GIF to image conversions, ensuring animated GIFs are not processed
def gif_to_image(input_path: str, output_path: str, target_format: str) -> str:
    fmt = target_format.upper()
    pil_format = _get_pil_format(fmt)

    with Image.open(input_path) as img:
        # Detect animated GIF
        try:
            n_frames = getattr(img, "n_frames", 1)
        except Exception:
            n_frames = 1

        if n_frames > 1:
            raise ValueError(
                "Animated GIFs are not supported in v1.1.0. "
                "Please convert a static GIF file."
            )

        # Convert palette mode to RGB/RGBA
        if img.mode == "P":
            img = img.convert("RGBA")
        if pil_format in ("JPEG",) and img.mode in ("RGBA", "LA"):
            img = img.convert("RGB")

        if pil_format == "JPEG":
            img.save(output_path, format=pil_format, quality=DEFAULT_IMAGE_QUALITY)
        else:
            img.save(output_path, format=pil_format)

    return output_path

def ico_to_png(input_path: str, output_path: str) -> str:
    with Image.open(input_path) as img:
        # ICO files contain multiple sizes — pick the largest
        sizes = getattr(img, "ico", None)
        if sizes:
            try:
                img.size = max(img.ico.sizes())
            except Exception:
                pass
        if img.mode != "RGBA":
            img = img.convert("RGBA")
        img.save(output_path, format="PNG")
    return output_path


def image_to_ico(input_path: str, output_path: str) -> str:
    ico_sizes = [(16, 16), (32, 32), (48, 48), (64, 64), (128, 128), (256, 256)]

    with Image.open(input_path) as img:
        # Convert to RGBA for transparency support in ICO
        if img.mode != "RGBA":
            img = img.convert("RGBA")

        w, h = img.size

        # Pad non-square images to square canvas preserving aspect ratio
        if w != h:
            size = max(w, h)
            canvas = Image.new("RGBA", (size, size), (255, 255, 255, 0))
            offset = ((size - w) // 2, (size - h) // 2)
            canvas.paste(img, offset)
            img = canvas

        # Resize to 256x256 as base, Pillow embeds all sizes automatically
        img = img.resize((256, 256), Image.LANCZOS)
        img.save(output_path, format="ICO", sizes=ico_sizes)

    return output_path


def image_to_pdf(input_path: str, output_path: str) -> str:
    with Image.open(input_path) as img:
        img = _prepare_image(img, "PDF")
        img.save(output_path, format="PDF", resolution=DEFAULT_IMAGE_QUALITY)
    return output_path


def image_to_pptx(input_path: str, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(input_path, Inches(0), Inches(0), Inches(10), Inches(7.5))
    prs.save(output_path)
    return output_path


def image_to_docx(input_path: str, output_path: str) -> str:
    from docx import Document
    from docx.shared import Inches as DocxInches
    doc = Document()
    doc.add_picture(input_path, width=DocxInches(6))
    doc.save(output_path)
    return output_path